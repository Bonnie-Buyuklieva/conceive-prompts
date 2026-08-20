#!/usr/bin/env python3
"""
Append a poster slide for a prompt to the deck. Append-only by design.

Usage:
    python poster/make_slide.py prompts/claims_v7.txt poster/poster_deck.pptx
        [--badge "PROMPT v7 · claims · CANDIDATE"] [--font-size 8] [--out other.pptx]

The deck's LAST slide must be the TEMPLATE: a slide with shapes named
'prompt_body' and 'version_badge', the badge text exactly 'TEMPLATE'.
The script deep-copies it, inserts the copy just before the template, fills
'prompt_body' with the .txt content (Consolas, one paragraph per line) and
sets 'version_badge'. Existing slides — including all manual annotation
cards, dot anchors and leader lines — are never touched.

Badge defaults to "PROMPT <version> · <task> · <date>" parsed from the
filename (<task>_<version>.txt). Font size defaults to 10 pt, dropping to
8 pt for prompts longer than 90 lines.
"""
from __future__ import annotations

import argparse
import datetime
import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

PROMPT_FONT = "Consolas"
PROMPT_COLOR = RGBColor(0x14, 0x14, 0x14)
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def find_template(prs):
    for idx, slide in enumerate(prs.slides):
        names = {sh.name: sh for sh in slide.shapes}
        if "prompt_body" in names and "version_badge" in names:
            if names["version_badge"].text_frame.text.strip().upper().startswith("TEMPLATE"):
                return idx, slide
    sys.exit("No TEMPLATE slide found (needs shapes 'prompt_body' + 'version_badge' text starting with TEMPLATE).")


def duplicate_slide(prs, src_slide):
    """Copy src_slide's shape tree onto a new blank slide (native shapes only)."""
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    spTree = new_slide.shapes._spTree
    # drop placeholders the blank layout brought in
    for sp in list(spTree):
        if sp.tag.endswith("}sp") or sp.tag.endswith("}pic") or sp.tag.endswith("}grpSp") \
           or sp.tag.endswith("}graphicFrame") or sp.tag.endswith("}cxnSp"):
            spTree.remove(sp)
    src_tree = src_slide.shapes._spTree
    for sp in src_tree:
        if sp.tag.endswith("}nvGrpSpPr") or sp.tag.endswith("}grpSpPr"):
            continue
        spTree.append(deepcopy(sp))
    return new_slide


def move_before(prs, moving_idx, anchor_idx):
    """Reorder sldIdLst so slide at moving_idx sits just before anchor_idx."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    moving = ids[moving_idx]
    lst.remove(moving)
    lst.insert(anchor_idx, moving)


def set_text(shape, text):
    """Replace text preserving the first run's formatting."""
    tf = shape.text_frame
    first_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
    font_kwargs = {}
    if first_run is not None:
        f = first_run.font
        font_kwargs = dict(name=f.name, size=f.size, bold=f.bold, italic=f.italic)
        color = f.color.rgb if f.color and f.color.type is not None else None
    else:
        color = None
    # wipe extra paragraphs
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text
    if font_kwargs.get("name"):
        run.font.name = font_kwargs["name"]
    if font_kwargs.get("size"):
        run.font.size = font_kwargs["size"]
    run.font.bold = font_kwargs.get("bold")
    run.font.italic = font_kwargs.get("italic")
    if color is not None:
        run.font.color.rgb = color


def fill_prompt_body(shape, lines, font_size_pt):
    tf = shape.text_frame
    tf.word_wrap = True
    # clear everything
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

    def style(run):
        run.font.name = PROMPT_FONT
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = PROMPT_COLOR

    first = True
    for line in lines:
        para = p0 if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text = line if line else " "
        style(run)
        pPr = para._p.get_or_add_pPr()
        # exact single spacing so the column height is predictable
        for tag in ("lnSpc",):
            for el in pPr.findall(f"{{{NS_A}}}{tag}"):
                pPr.remove(el)


def derive_badge(prompt_path: Path) -> str:
    """'PROMPT <version> · <task> · <date>' from a <task>_<version>.txt filename."""
    m = re.fullmatch(r"(?P<task>.+?)_(?P<ver>v\d+|A\d+)", prompt_path.stem, re.IGNORECASE)
    if m:
        task, version = m.group("task"), m.group("ver")
    else:
        task, version = "prompt", prompt_path.stem
    stamp = datetime.date.today().isoformat()
    return f"PROMPT {version} · {task} · {stamp}"


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("prompt_txt")
    ap.add_argument("deck")
    ap.add_argument("--badge", default=None,
                    help="badge text (default: derived from the filename)")
    ap.add_argument("--font-size", type=float, default=None,
                    help="prompt body pt size (default: 10, or 8 for >90 lines)")
    ap.add_argument("--out", default=None, help="output deck (default: in place)")
    args = ap.parse_args()

    prompt_path = Path(args.prompt_txt)
    lines = prompt_path.read_text(encoding="utf-8").splitlines()

    badge = args.badge or derive_badge(prompt_path)
    font_size = args.font_size or (10 if len(lines) <= 90 else 8)

    prs = Presentation(args.deck)
    tmpl_idx, tmpl = find_template(prs)

    new_slide = duplicate_slide(prs, tmpl)
    move_before(prs, len(prs.slides._sldIdLst) - 1, tmpl_idx)

    names = {sh.name: sh for sh in new_slide.shapes}
    fill_prompt_body(names["prompt_body"], lines, font_size)
    set_text(names["version_badge"], badge)

    out = args.out or args.deck
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"Appended slide {tmpl_idx + 1}/{n} from {prompt_path.name} -> {out}")
    print(f"  badge: {badge!r} | font: {PROMPT_FONT} {font_size}pt | lines: {len(lines)}")
    print("  Template remains the last slide. Existing slides untouched.")


if __name__ == "__main__":
    main()
